import os
import sys
import re
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# --- 通用視覺參數設定 ---
# 使用中性的深藍灰色作為主題色
THEME_COLOR = RGBColor(44, 62, 80) 
TEXT_DARK = RGBColor(33, 33, 33)
FONT_MAIN = 'Microsoft JhengHei'

def clean_md_syntax(text):
    """徹底清除 Markdown 語法符號與特定報告標籤"""
    if not text: return ""
    # 移除粗體/斜體語法
    text = re.sub(r'[\*_]{1,2}(.*?)[\*_]{1,2}', r'\1', text)
    # 移除條列符號與多餘空格
    text = text.replace('*', '').strip()
    # 移除常見的引導標籤，使內容更簡潔
    labels_to_remove = ['協定名稱：', '技術名稱：', '任務：']
    for label in labels_to_remove:
        text = text.replace(label, '')
    return text

def set_title_style(title_shape, main_title, sub_title, page_idx=0, total_pages=1):
    """動態組合標題並加入分頁標示，移除寫死的品牌字眼"""
    full_title = clean_md_syntax(main_title)
    if sub_title:
        full_title += f" - {clean_md_syntax(sub_title)}"
    
    # 加入分頁進度標記 (例如: 1/2)
    if total_pages > 1:
        full_title += f" ({page_idx + 1}/{total_pages})"

    title_shape.text = full_title
    tf = title_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.name = FONT_MAIN
    p.font.bold = True
    p.font.color.rgb = THEME_COLOR
    
    # 標題大小根據長度自適應
    length = len(full_title)
    if length > 35: p.font.size = Pt(18)
    elif length > 25: p.font.size = Pt(22)
    else: p.font.size = Pt(26)

def get_body_placeholder(slide):
    """取得投影片中的內容預留位置"""
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 2 or shape.placeholder_format.idx == 1:
            return shape
    return None

def create_pptx_from_md(md_file="report.md", output_file="Research_Presentation.pptx"):
    if not os.path.exists(md_file): 
        print(f"❌ 找不到檔案: {md_file}")
        return
    
    # 嘗試從檔名解析主題，作為首頁標題
    base_name = os.path.basename(md_file).replace('.md', '')
    display_subject = base_name.split('_')[-1] if '_' in base_name else "技術研究"

    prs = Presentation()
    with open(md_file, "r", encoding="utf-8") as f:
        content = f.read()

    # --- 1. 處理封面 (通用標題) ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"{display_subject}\n技術分析與研究報告"
    slide.placeholders[1].text = f"AI 輔助生成文件\n日期：{time.strftime('%Y-%m-%d')}"

    # --- 2. 解析大章節 (##) ---
    main_sections = re.split(r'\n(?=## )', content)
    
    for m_sec in main_sections:
        m_lines = m_sec.strip().split('\n')
        if not m_lines: continue
        
        main_title = m_lines[0].strip('# ')
        
        # --- 3. 解析小章節 (###) ---
        sub_sections = re.split(r'\n(?=### )', '\n'.join(m_lines[1:]))
        
        for s_sec in sub_sections:
            s_lines = s_sec.strip().split('\n')
            if not s_lines: continue
            
            if s_lines[0].startswith('###'):
                sub_title = s_lines[0].strip('# ')
                body_lines = s_lines[1:]
            else:
                sub_title = ""
                body_lines = s_lines

            # 分離純文字內容與表格數據
            table_data = []
            text_content = []
            for line in body_lines:
                line = line.strip()
                if not line or line.startswith('---'): continue
                if line.startswith('|'):
                    if '---' in line: continue
                    cells = [c.strip() for c in line.split('|') if c.strip()]
                    if cells: table_data.append(cells)
                else:
                    text_content.append(line)

            # 🚀 處理文字分頁 (優化排版)
            if text_content:
                chunk_size = 7 # 每頁最多行數
                pages = [text_content[i:i + chunk_size] for i in range(0, len(text_content), chunk_size)]
                for idx, chunk in enumerate(pages):
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    set_title_style(slide.shapes.title, main_title, sub_title, idx, len(pages))
                    
                    body = get_body_placeholder(slide)
                    if body:
                        tf = body.text_frame
                        tf.clear()
                        for text_line in chunk:
                            p = tf.add_paragraph()
                            p.text = clean_md_syntax(text_line)
                            p.font.name = FONT_MAIN
                            p.font.size = Pt(16)
                            p.space_after = Pt(10)

            # 🚀 處理表格分頁
            if table_data:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                set_title_style(slide.shapes.title, main_title, sub_title + " (詳細數據)")
                
                body = get_body_placeholder(slide)
                if body:
                    sp = body.element
                    sp.getparent().remove(sp) # 移除佔位框改放表格

                rows, cols = len(table_data), len(table_data[0])
                table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9.0), Inches(0.4))
                table = table_shape.table
                for r, row in enumerate(table_data):
                    for c, cell_text in enumerate(row):
                        cell = table.cell(r, c)
                        cell.text = clean_md_syntax(cell_text)
                        for p in cell.text_frame.paragraphs:
                            p.font.name = FONT_MAIN
                            p.font.size = Pt(10)
                            if r == 0: # 設定通用表頭樣式
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = THEME_COLOR
                                p.font.color.rgb = RGBColor(255, 255, 255)
                                p.font.bold = True

    prs.save(output_file)
    print(f"✅ 通用版 PPT 已生成至: {output_file}")

if __name__ == "__main__":
    # 支援動態指定輸入檔案，預設為 generic_report.md
    target = sys.argv[1] if len(sys.argv) > 1 else "report.md"
    create_pptx_from_md(target)
