import os
import sys
import time
from google import genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_pptx_from_text(text, filename="Vector_Research.pptx"):
    """將 AI 輸出的結構化文字轉為美化後的 PPTX 檔案"""
    prs = Presentation()
    
    # --- 1. 加入總標題首頁 ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Vector Informatik 技術研究報告"
    subtitle.text = f"由 AI 助理自動生成\n日期：{time.strftime('%Y-%m-%d')}"
    
    # 設定首頁字體為微軟正黑體
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.name = 'Microsoft JhengHei'
                paragraph.font.bold = True

    # --- 2. 處理內容投影片 ---
    sections = text.split('Slide:')
    for section in sections:
        if not section.strip():
            continue
        
        lines = section.strip().split('\n')
        title_text = lines[0].strip()
        bullet_points = lines[1:]
        
        # 加入一張「標題+內容」的投影片 (版面 1)
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # --- 優化標題格式 ---
        title_shape = slide.shapes.title
        title_shape.text = title_text
        title_frame = title_shape.text_frame
        for p in title_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(32)
            p.font.name = 'Microsoft JhengHei'
            p.font.color.rgb = RGBColor(0, 51, 102) # Vector 深藍色風格

        # --- 優化內文內容 ---
        if bullet_points:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.word_wrap = True # 自動換行
            
            for point in bullet_points:
                # 移除 AI 產出的多餘符號如 * 或 -
                clean_point = point.strip().lstrip('-').lstrip('*').strip()
                if clean_point:
                    p = tf.add_paragraph()
                    p.text = clean_point
                    p.font.size = Pt(20)
                    p.font.name = 'Microsoft JhengHei'
                    p.space_after = Pt(12) # 設定段落間距，避免文字太擠
                    p.level = 0
    
    prs.save(filename)
    print(f"✅ 成功產生優化美化版 PPT 檔案：{filename}", file=sys.stderr)

def main():
    api_key = os.environ.get("GEMINI_API_KEY")
    issue_text = os.environ.get("ISSUE_BODY", "沒有研究主題")
    label = os.environ.get("TRIGGER_LABEL", "research")
    
    if not api_key:
        print("❌ 錯誤：找不到 API KEY")
        return

    client = genai.Client(api_key=api_key)
    
    # 🧠 根據標籤切換 Prompt 與身份
    if label == "presentation":
        system_instruction = "你是一位資深的車用電子技術顧問，擅長製作 Vector Informatik 產品方案的專業簡報。"
        prompt = f"""
        請針對以下主題製作 PPT 結構，直接以文字格式輸出。
        
        【內容架構要求】：
        1. 議程 (Agenda)：列出本次研究的關鍵章節。
        2. 充電協定解析：主動找出對應的 ISO 15118 (-2, -20), DIN 70121 等標準。
        3. 硬體工具映射：對應 VH5110A, VN 系列, VT System 等 Vector 設備。
        4. 技術挑戰分析：深入分析實作中的技術門檻（如 PLC 訊號、TLS 加密）。
        5. 總結與展望。

        【格式規範】：
        1. 每張投影片以 'Slide:' 開頭。
        2. 內容精煉，每張投影片不超過 5 個重點。
        3. 每點字數建議在 15-20 字以內。
        
        Slide: [投影片標題]
        - [重點 1]
        - [重點 2]

        主題：{issue_text}
        """
    else:
        system_instruction = "你是一位專精於 Vector Informatik 與 EV 充電標準的技術研究專家。"
        prompt = f"請針對以下主題產出詳細的 Markdown 技術研究報告，包含通訊協定、硬體對應與技術分析：\n\n{issue_text}"

    max_attempts = 3
    base_delay = 20 

    for i in range(max_attempts):
        try:
            print(f"📡 第 {i+1} 次嘗試 ({label} 模式)...", file=sys.stderr)
            response = client.models.generate_content(
                model='gemini-2.5-flash', 
                contents=f"{system_instruction}\n\n{prompt}"
            )
            
            if response.text:
                # 輸出文字結果 (給 Issue 留言用)
                print(response.text)
                
                # 如果是簡報標籤，則執行 PPT 製作
                if label == "presentation":
                    create_pptx_from_text(response.text)
                return 
            
        except Exception as e:
            error_msg = str(e)
            if i < max_attempts - 1:
                wait_time = base_delay * (2 ** i)
                print(f"⚠️ 嘗試失敗，將於 {wait_time} 秒後重試...", file=sys.stderr)
                time.sleep(wait_time)
            else:
                print(f"### ❌ AI 員工在 3 次重試後仍然失敗\n最後錯誤訊息：`{error_msg}`")

if __name__ == "__main__":
    main()
